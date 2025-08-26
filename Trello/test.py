from trello import TrelloClient
from pptx import Presentation
from copy import deepcopy
import time

start_time = time.time()

#----------Configuração do Trello----------
client = TrelloClient(
    api_key='a2ffb6570256b8afa8d0262596294e36',
    token='ATTA71a7e45fb51ca2bd4c24dd8d44a032c98a778e0ff331eb6f70130e3ec24e8a31AAFA42DF'
)

#----------Criar listas-------------
board_ID_list = []
list_ID_list = []
card_ID_list = []

#----------Acessando as listas----------
board = client.get_board('66ddcdc0176ed00c5c401bb5')
pendente = client.get_list('66ddcdc0176ed00c5c401bbc')
andamento = client.get_list('66ddcdc0176ed00c5c401bbd')
concluido = client.get_list('66ddcdc0176ed00c5c401bbe')
correcao = client.get_list('66ddd0f27a2d8828175760d8')

#-----------Functions-------------------

def get_cards(membro_id, list):
    list = client.get_list(list)
    cards = list.list_cards()
    lista_cards = []
    for card in cards:
        if membro_id in card.member_id:
            lista_cards.append(card)
    return lista_cards


#-----------Lista de membros-------------

members = board.all_members()
list_membros = []

for member in members:
    list_membros.append(member)  


#--------Criar dicionário com usuários------

dados_membros = {} #armazena subdicionários
total = 0
concluido = 0
andamento = 0
pendente = 0
correcao = 0


for membro in list_membros:
    crd_pendente = get_cards(membro.id, '66ddcdc0176ed00c5c401bbc')
    crd_andamento = get_cards(membro.id,'66ddcdc0176ed00c5c401bbd') 
    crd_concluido = get_cards(membro.id,'66ddcdc0176ed00c5c401bbe')
    crd_correcao = get_cards(membro.id,'66ddd0f27a2d8828175760d8')

    dados_membros[membro] = {
        "pendente": crd_pendente,
        "andamento": crd_andamento,
        "concluido": crd_concluido,
        "correção": crd_correcao,
        "total": sum(len(lst) for lst in [crd_pendente, crd_andamento, crd_concluido, crd_correcao])
    }

    pendente += len(crd_pendente)
    andamento += len(crd_andamento)
    concluido += len(crd_concluido)
    correcao += len(crd_correcao)
    total += dados_membros[membro]["total"]


#----------Imprimir dados----------------
print("Slide 1: Overview")
print(f"Total de cards: {total}")
print(f"Concluídos: {concluido}")#Concluídos
print(f"Em andamento: {andamento}")#Em andamento
print(f"Pendentes: {pendente}")#Pendentes


#----------Imprimir dados de membro específico----------
for i, membro in enumerate(list_membros, start=2):
    print("\n" + "=" * 20)
    print(f"Slide {i}: {membro.full_name}")

    # Get card titles for each category
    pendente = "| ".join(card.name for card in dados_membros[membro]["pendente"])
    andamento = "| ".join(card.name for card in dados_membros[membro]["andamento"])
    concluido = "| ".join(card.name for card in dados_membros[membro]["concluido"])
    correcao = "| ".join(card.name for card in dados_membros[membro]["correção"])

    print(f"Concluídos: {concluido}")
    print(f"Andamento: {andamento}")
    print(f"Pendentes: {pendente}")
    print(f"Correções: {correcao}")
    print(f"Total de cards: {dados_membros[membro]['total']}")


print("Slide 3: Gráfico? Será que dá?")

#Criar a apresentação
presentation = Presentation("tutorial_presentation 1.pptx")


# Slide de overview (índice 1)
overview_slide = presentation.slides[1]

# Nomes reais das listas
nomes_listas = {
    "concluido": "Concluído",
    "andamento": "Em andamento",
    "pendente": "Pendente",
    "correção": "Correção"
}

# Totais por categoria já calculados
totais_por_categoria = {
    "concluido": sum(len(dados_membros[m]["concluido"]) for m in list_membros),
    "andamento": sum(len(dados_membros[m]["andamento"]) for m in list_membros),
    "pendente": sum(len(dados_membros[m]["pendente"]) for m in list_membros),
    "correção": sum(len(dados_membros[m]["correção"]) for m in list_membros)
}

# Substituir os placeholders no slide
for shape in overview_slide.shapes:
    if shape.has_text_frame:
        texto = shape.text_frame.text
        for i, cat in enumerate(nomes_listas, start=1):
            texto = texto.replace(f"List tittle {i}", nomes_listas[cat])
            texto = texto.replace(f"Card {i} quantity", str(totais_por_categoria[cat]))
        shape.text_frame.text = texto

# Remover apenas placeholders não substituídos do overview
for shape in list(overview_slide.shapes):
    if shape.has_text_frame:
        texto = shape.text_frame.text.strip()
        if texto.startswith("List tittle") or texto.startswith("Card") and "quantity" in texto:
            shape._element.getparent().remove(shape._element)



# Função para adicionar slides
def copiar_slide(prs, slide_index):
    slide_origem = prs.slides[slide_index]
    layout = slide_origem.slide_layout
    novo_slide = prs.slides.add_slide(layout)

    for shape in slide_origem.shapes:
        novo_shape = deepcopy(shape.element)
        novo_slide.shapes._spTree.insert_element_before(novo_shape, 'p:extLst')

    return presentation.slides[-1]  # Retorna o novo slide criado
        

#Slide for user
for membro in list_membros:

    slide = copiar_slide(presentation, 2)

    #identificar placeholder de "User Name" e substituir pelo nome do membro

    for shape in slide.shapes:
        if "User Name" in shape.text_frame.text:
            shape.text_frame.text = shape.text_frame.text.replace("User Name", membro.full_name)


    for shape in slide.shapes:
        if shape.has_text_frame:
            # Junta os cards atribuídos ao membro atual com suas categorias
            cards_com_categoria = []
            for categoria in ["concluido", "andamento", "pendente", "correção"]:
                for card in dados_membros[membro][categoria]:
                    if membro.id in getattr(card, "member_ids", []):
                        cards_com_categoria.append((categoria.capitalize(), card.name))

            # Substitui até 14 pares válidos
            for i in range(1, 15):
                placeholder_categoria = f"Card list {i}"
                placeholder_nome = f"Card {i} name"

                if i <= len(cards_com_categoria):
                    categoria_nome, card_nome = cards_com_categoria[i - 1]
                else:
                    categoria_nome, card_nome = "", ""

                # percorre cada parágrafo e substitui apenas o que corresponde
                for p in shape.text_frame.paragraphs:
                    if placeholder_categoria in p.text:
                        p.text = p.text.replace(placeholder_categoria, categoria_nome)
                    if placeholder_nome in p.text:
                        p.text = p.text.replace(placeholder_nome, card_nome)


# Limpar textos indesejados como "Concluido0", "Andamento1", etc.
def remover_textos_indesejados(presentation):
    categorias = ["Concluido", "Andamento", "Pendente", "Correção"]
    for slide in presentation.slides:
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                texto = shape.text_frame.text.strip()
                for cat in categorias:
                    if texto.startswith(cat) and texto[len(cat):].isdigit():
                        shape._element.getparent().remove(shape._element)
                        break

def remover_caixas_vazias(presentation):
    for slide in presentation.slides:
        for shape in list(slide.shapes):
            if shape.has_text_frame and not shape.text_frame.text.strip():
                shape._element.getparent().remove(shape._element)

# Chamar as funções
remover_textos_indesejados(presentation)
remover_caixas_vazias(presentation)







#categorias = ["Concluido", "Andamento", "Pendente", "Correção"]
#for slide in presentation.slides:
#    for shape in list(slide.shapes):
#        if shape.has_text_frame:
#            texto = shape.text_frame.text.strip()
#            # Se texto começar com qualquer categoria e terminar em número, remove
#            for cat in categorias:
#                if texto.startswith(cat) and texto[len(cat):].isdigit():
#                    shape._element.getparent().remove(shape._element)
#                    break
#
#
##Excluir caixas vazias
#for slide in presentation.slides: #Exclui caixas vazias
#    for shape in list(slide.shapes):  # list() evita erro ao remover durante iteração
#        if shape.has_text_frame and not shape.text_frame.text.strip():
#            shape._element.getparent().remove(shape._element)
#
presentation.save("tutorial_presentation.pptx")

end_time = time.time()

print(f'tempo de execução: {end_time - start_time:.2f} segundos')