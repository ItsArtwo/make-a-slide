# /Trello/trello_report/slide_generator.py
from pptx import Presentation
from copy import deepcopy
def copiar_slide(prs, slide_index):
    slide_origem = prs.slides[slide_index]
    layout = slide_origem.slide_layout
    novo_slide = prs.slides.add_slide(layout)
    for shape in slide_origem.shapes:
        novo_shape = deepcopy(shape.element)
        novo_slide.shapes._spTree.insert_element_before(novo_shape, 'p:extLst')
    return novo_slide
def excluir_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])
def limpar_placeholders_e_textos_vazios(prs):
    """
    Remove placeholders não substituídos, textos como 'Concluído0', 'Pendente1',
    e caixas de texto vazias da apresentação.
    """
    # Remover placeholders não substituídos do slide de overview
    overview_slide = prs.slides[1]
    for shape in list(overview_slide.shapes):
        if shape.has_text_frame:
            texto = shape.text_frame.text.strip()
            if texto.startswith("List tittle") or (texto.startswith("Card") and "quantity" in texto):
                shape._element.getparent().remove(shape._element)
    # Remover textos "Concluído0, Concluído1..." e outros placeholders numéricos
    categorias = ["Concluido", "Concluído", "Andamento", "Pendente", "Correção"]
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                texto = shape.text_frame.text.strip()
                for cat in categorias:
                    if texto.startswith(cat) and texto[len(cat):].isdigit():
                        shape._element.getparent().remove(shape._element)
                        break
    # Excluir caixas de texto vazias
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.has_text_frame and not shape.text_frame.text.strip():
                shape._element.getparent().remove(shape._element)
def gerar_apresentacao(membros, dados_membros, totais):
    prs = Presentation("tutorial_presentation 1.pptx")
    # Atualiza o slide de overview
    overview_slide = prs.slides[1]
    nomes_listas = {
        "concluido": "Concluído",
        "andamento": "Em andamento",
        "pendente": "Pendente",
        "correcao": "Correção"
    }
    for shape in overview_slide.shapes:
        if shape.has_text_frame:
            texto = shape.text_frame.text
            for i, cat in enumerate(nomes_listas, start=1):
                texto = texto.replace(f"List tittle {i}", nomes_listas[cat])
                texto = texto.replace(f"Card {i} quantity", str(totais[cat]))
            shape.text_frame.text = texto
    # Criar slides por membro
    for membro in membros:
        slide = copiar_slide(prs, 2)
        for shape in slide.shapes:
            if shape.has_text_frame and "User Name" in shape.text_frame.text:
                shape.text_frame.text = shape.text_frame.text.replace("User Name", membro.full_name)
            # Substituir placeholders de cards
            cards_com_categoria = []
            for categoria in ["concluido", "andamento", "pendente", "correcao"]:
                for card in dados_membros[membro][categoria]:
                    cards_com_categoria.append((categoria.capitalize(), card.name))
            for i in range(1, 15): # Iterar até 14 para cobrir Card list 1 a Card list 14
                placeholder_categoria = f"Card list {i}"
                placeholder_nome = f"Card {i} name"
                categoria_nome, card_nome = "", ""
                if i <= len(cards_com_categoria):
                    categoria_nome, card_nome = cards_com_categoria[i - 1]
                # Percorre cada parágrafo e substitui apenas o que corresponde
                for p in shape.text_frame.paragraphs:
                    if placeholder_categoria in p.text:
                        p.text = p.text.replace(placeholder_categoria, categoria_nome)
                    if placeholder_nome in p.text:
                        p.text = p.text.replace(placeholder_nome, card_nome)
    # Excluir o slide 3 (índice 2) que é o template para os slides de membro
    excluir_slide(prs, 2)
    # Chamar a nova função de limpeza
    limpar_placeholders_e_textos_vazios(prs)
    prs.save("tutorial_presentation.pptx")