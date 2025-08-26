from trello import TrelloClient
from pptx import Presentation
import time

start_time = time.time()

#----------Configuração do Trello----------
client = TrelloClient(
    api_key='a2ffb6570256b8afa8d0262596294e36',
    token='ATTA71a7e45fb51ca2bd4c24dd8d44a032c98a778e0ff331eb6f70130e3ec24e8a31AAFA42DF'
)

#----------Criar listas-------------
board_ID_list = []
list_ID_list = []
card_ID_list = []

#----------Acessando as listas----------
board = client.get_board('66ddcdc0176ed00c5c401bb5')
pendente = client.get_list('66ddcdc0176ed00c5c401bbc')
andamento = client.get_list('66ddcdc0176ed00c5c401bbd')
concluido = client.get_list('66ddcdc0176ed00c5c401bbe')
correcao = client.get_list('66ddd0f27a2d8828175760d8')

#-----------Functions-------------------

def get_cards(membro_id, list):
    list = client.get_list(list)
    cards = list.list_cards()
    lista_cards = []
    for card in cards:
        if membro_id in card.member_id:
            lista_cards.append(card)
    return lista_cards


#-----------Lista de membros-------------

members = board.all_members()
list_membros = []

for member in members:
    list_membros.append(member)  


#--------Criar dicionário com usuários------

dados_membros = {} #armazena subdicionários
total = 0
concluido = 0
andamento = 0
pendente = 0
correcao = 0


for membro in list_membros:
    crd_pendente = get_cards(membro.id, '66ddcdc0176ed00c5c401bbc')
    crd_andamento = get_cards(membro.id,'66ddcdc0176ed00c5c401bbd') 
    crd_concluido = get_cards(membro.id,'66ddcdc0176ed00c5c401bbe')
    crd_correcao = get_cards(membro.id,'66ddd0f27a2d8828175760d8')

    dados_membros[membro] = {
        "pendente": crd_pendente,
        "andamento": crd_andamento,
        "concluido": crd_concluido,
        "correção": crd_correcao,
        "total": sum(len(lst) for lst in [crd_pendente, crd_andamento, crd_concluido, crd_correcao])
    }

    pendente += len(crd_pendente)
    andamento += len(crd_andamento)
    concluido += len(crd_concluido)
    correcao += len(crd_correcao)
    total += dados_membros[membro]["total"]


#----------Imprimir dados----------------
print("Slide 1: Overview")
print(f"Total de cards: {total}")
print(f"Concluídos: {concluido}")#Concluídos
print(f"Em andamento: {andamento}")#Em andamento
print(f"Pendentes: {pendente}")#Pendentes


#----------Imprimir dados de membro específico----------
for i, membro in enumerate(list_membros, start=2):
    print("\n" + "=" * 20)
    print(f"Slide {i}: {membro.full_name}")

    # Get card titles for each category
    pendente = "| ".join(card.name for card in dados_membros[membro]["pendente"])
    andamento = "| ".join(card.name for card in dados_membros[membro]["andamento"])
    concluido = "| ".join(card.name for card in dados_membros[membro]["concluido"])
    correcao = "| ".join(card.name for card in dados_membros[membro]["correção"])

    print(f"Concluídos: {concluido}")
    print(f"Andamento: {andamento}")
    print(f"Pendentes: {pendente}")
    print(f"Correções: {correcao}")
    print(f"Total de cards: {dados_membros[membro]['total']}")


print("Slide 3: Gráfico? Será que dá?")

#Criar a apresentação
presentation = Presentation("tutorial_presentation 1.pptx")

#Adicionar slide

title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Overview"
slide.placeholders[1].text = "Primeiríssima reunion"

#Primeiro slide
overview_slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(overview_slide_layout)
slide.shapes.title.text = 'Overview'
slide.placeholders[1].text = (
    f"Total de cards: {total}\n"
    f"Concluídos: {concluido}\n"
    f"Em andamento: {andamento}\n"
    f"Pendentes: {pendente}"
)

#Slide for user

for membro in list_membros:
    user_slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(user_slide_layout)
    slide.shapes.title.text = f"{membro.full_name}"
    slide.placeholders[1].text = (
        f"Concluídos: {' | '.join(card.name for card in dados_membros[membro]['concluido'])}\n"
        f"Andamento: {' | '.join(card.name for card in dados_membros[membro]['andamento'])}\n"
        f"Pendentes: {' | '.join(card.name for card in dados_membros[membro]['pendente'])}\n"
        f"Correções: {' | '.join(card.name for card in dados_membros[membro]['correção'])}\n"
        f"Total de cards: {dados_membros[membro]['total']}"
    )

presentation.save("tutorial_presentation.pptx")

end_time = time.time()

print(f'tempo de execução: {end_time - start_time:.2f} segundos')